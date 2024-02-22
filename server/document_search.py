from llama_index.core import Document, GPTVectorStoreIndex
import fitz
from pptx import Presentation
from docx import Document as DocxDocument
from openpyxl import load_workbook
import os

def create_document_from_word(docx_path):
    document = DocxDocument(docx_path)
    text = ""
    for paragraph in document.paragraphs:
        text += paragraph.text + '\n'
    return Document(text=text)

def create_document_from_excel(xlsx_path):
    workbook = load_workbook(xlsx_path)
    text = ""
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            text += ' '.join(str(cell) for cell in row) + '\n'
    return Document(text=text)

def create_document_from_pdf(pdf_path):
    with fitz.open(pdf_path) as pdf_document:
        text = ""
        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            text += page.get_text()

    return Document(text=text)


def create_document_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text

    return Document(text=text)


def create_document_from_file(file_path):
    if file_path.endswith(".pdf"):
        return create_document_from_pdf(file_path)
    elif file_path.endswith(".pptx") or file_path.endswith(".ppt"):
        return create_document_from_pptx(file_path)
    elif file_path.endswith(".docx") or file_path.endswith(".doc"):
        return create_document_from_word(file_path)
    elif file_path.endswith(".xlsx") or file_path.endswith(".xls"):
        return create_document_from_excel(file_path)
    else:
        return None


def get_files_in_directory(directory_path):
    return [os.path.join(directory_path, filename) for filename in os.listdir(directory_path)]

def index_documents(search):
    documents = []
    for file_path in get_files_in_directory("./data"):
        document = create_document_from_file(file_path)
        if document:
            documents.append(document)
    index = GPTVectorStoreIndex.from_documents(documents)
    query_engine = index.as_query_engine()
    response = query_engine.query(search)
    return response


